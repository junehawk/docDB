"""
Microsoft Office 형식 추출기들.
DOCX, PPTX, XLS, XLSX 파일 추출을 담당합니다.
"""

from loguru import logger

from .base_extractor import BaseExtractor, ExtractionResult


class DocxExtractor(BaseExtractor):
    """
    DOCX (Word) 파일 추출기.
    python-docx 라이브러리를 사용하여 문단과 표를 추출합니다.
    """

    def extract(self) -> ExtractionResult:
        """
        DOCX 파일에서 텍스트를 추출합니다.

        Returns:
            ExtractionResult: 추출 결과
        """
        if not self.validate_file():
            return self._create_error_result("File validation failed")

        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed")
            return self._create_error_result("python-docx not installed")

        try:
            doc = Document(str(self.file_path))
            text_parts = []

            # 문서 내장 properties 추출
            doc_properties = {}
            try:
                props = doc.core_properties
                doc_properties = {
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'keywords': props.keywords or '',
                    'created': props.created.isoformat() if props.created else '',
                    'modified': props.modified.isoformat() if props.modified else '',
                    'category': props.category or '',
                }
            except Exception:
                doc_properties = {}

            # 문단 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 표 추출
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    text_parts.append(f"\n[Table]\n{table_text}\n")

            if text_parts:
                text = "\n".join(text_parts)
                return self._create_success_result(
                    text, {
                        "method": "python-docx",
                        "paragraphs": len(doc.paragraphs),
                        "tables": len(doc.tables),
                        "doc_properties": doc_properties,
                    }
                )
            else:
                return self._create_error_result("No text found in DOCX")

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return self._create_error_result(f"DOCX error: {str(e)}")

    @staticmethod
    def _extract_table_text(table) -> str:
        """
        표에서 텍스트를 추출합니다.

        Args:
            table: python-docx 테이블 객체

        Returns:
            str: 추출된 표 텍스트
        """
        table_rows = []

        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = "".join([p.text for p in cell.paragraphs]).strip()
                row_cells.append(cell_text)
            table_rows.append(" | ".join(row_cells))

        return "\n".join(table_rows)


class PptxExtractor(BaseExtractor):
    """
    PPTX (PowerPoint) 파일 추출기.
    python-pptx 라이브러리를 사용하여 슬라이드와 노트에서 텍스트를 추출합니다.
    """

    def extract(self) -> ExtractionResult:
        """
        PPTX 파일에서 텍스트를 추출합니다.

        Returns:
            ExtractionResult: 추출 결과
        """
        if not self.validate_file():
            return self._create_error_result("File validation failed")

        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed")
            return self._create_error_result("python-pptx not installed")

        try:
            prs = Presentation(str(self.file_path))
            text_parts = []

            # 문서 내장 properties 추출
            doc_properties = {}
            try:
                props = prs.core_properties
                doc_properties = {
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'keywords': props.keywords or '',
                    'created': props.created.isoformat() if props.created else '',
                    'modified': props.modified.isoformat() if props.modified else '',
                    'category': props.category or '',
                }
            except Exception:
                doc_properties = {}

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                # 슬라이드의 모든 도형에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)

                    # 테이블인 경우 처리
                    if shape.has_table:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_text.append(f"[Table]\n{table_text}")

                # 슬라이드 노트 추출
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text.append(f"[Notes]\n{notes_text}")

                if slide_text:
                    text_parts.append(f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text))

            if text_parts:
                text = "\n".join(text_parts)
                return self._create_success_result(
                    text, {
                        "method": "python-pptx",
                        "slides": len(prs.slides),
                        "doc_properties": doc_properties,
                    }
                )
            else:
                return self._create_error_result("No text found in PPTX")

        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return self._create_error_result(f"PPTX error: {str(e)}")

    @staticmethod
    def _extract_table_text(table) -> str:
        """
        표에서 텍스트를 추출합니다.

        Args:
            table: python-pptx 테이블 객체

        Returns:
            str: 추출된 표 텍스트
        """
        table_rows = []

        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_cells.append(cell_text)
            table_rows.append(" | ".join(row_cells))

        return "\n".join(table_rows)


class XlsExtractor(BaseExtractor):
    """
    XLS (구형 Excel) 파일 추출기.
    xlrd 라이브러리를 사용하여 OLE2 형식(.xls)의 시트 데이터를 추출합니다.
    """

    MAX_ROWS_PER_SHEET = 100

    def extract(self) -> ExtractionResult:
        if not self.validate_file():
            return self._create_error_result("File validation failed")

        try:
            import xlrd
        except ImportError:
            logger.error("xlrd not installed")
            return self._create_error_result("xlrd not installed. pip install xlrd")

        try:
            workbook = xlrd.open_workbook(str(self.file_path))
            text_parts = []
            sheets_processed = 0

            for sheet in workbook.sheets():
                if sheet.nrows == 0:
                    continue

                sheet_rows = []

                # 헤더 행
                headers = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
                if any(h.strip() for h in headers):
                    sheet_rows.append("Headers: " + ", ".join(headers))

                # 데이터 행
                max_row = min(sheet.nrows, self.MAX_ROWS_PER_SHEET + 1)
                for r in range(1, max_row):
                    row_values = []
                    for c in range(sheet.ncols):
                        val = sheet.cell_value(r, c)
                        if val not in (None, ''):
                            header = headers[c] if c < len(headers) else f'Col{c}'
                            row_values.append(f"{header}: {val}")
                    if row_values:
                        sheet_rows.append(", ".join(row_values))

                if sheet_rows:
                    text_parts.append(f"\n[Sheet: {sheet.name}]\n" + "\n".join(sheet_rows))
                    sheets_processed += 1

            if text_parts:
                text = "\n".join(text_parts)
                return self._create_success_result(
                    text, {
                        "method": "xlrd",
                        "sheets_processed": sheets_processed,
                        "total_sheets": workbook.nsheets,
                        "doc_properties": {},
                    }
                )
            else:
                return self._create_error_result("No data found in XLS")

        except Exception as e:
            logger.error(f"XLS extraction failed: {e}")
            return self._create_error_result(f"XLS error: {str(e)}")


class XlsxExtractor(BaseExtractor):
    """
    XLSX (Excel) 파일 추출기.
    openpyxl 라이브러리를 사용하여 시트 데이터를 추출합니다.
    시트당 첫 100행의 헤더와 셀 값을 추출합니다.
    """

    MAX_ROWS_PER_SHEET = 100  # 시트당 최대 행 수

    def extract(self) -> ExtractionResult:
        """
        XLSX 파일에서 텍스트를 추출합니다.

        Returns:
            ExtractionResult: 추출 결과
        """
        if not self.validate_file():
            return self._create_error_result("File validation failed")

        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.error("openpyxl not installed")
            return self._create_error_result("openpyxl not installed")

        try:
            workbook = load_workbook(str(self.file_path), data_only=True)
            text_parts = []
            sheets_processed = 0

            # 문서 내장 properties 추출
            doc_properties = {}
            try:
                props = workbook.properties
                doc_properties = {
                    'title': props.title or '',
                    'creator': props.creator or '',
                    'subject': props.subject or '',
                    'keywords': props.keywords or '',
                    'created': props.created.isoformat() if props.created else '',
                    'modified': props.modified.isoformat() if props.modified else '',
                }
            except Exception:
                doc_properties = {}

            for sheet_name in workbook.sheetnames:
                try:
                    ws = workbook[sheet_name]
                    sheet_text = self._extract_sheet_text(ws)

                    if sheet_text:
                        text_parts.append(f"\n[Sheet: {sheet_name}]\n{sheet_text}")
                        sheets_processed += 1
                except Exception as e:
                    logger.debug(f"Failed to process sheet '{sheet_name}': {e}")
                    continue

            if text_parts:
                text = "\n".join(text_parts)
                return self._create_success_result(
                    text, {
                        "method": "openpyxl",
                        "sheets_processed": sheets_processed,
                        "total_sheets": len(workbook.sheetnames),
                        "max_rows_per_sheet": self.MAX_ROWS_PER_SHEET,
                        "doc_properties": doc_properties,
                    }
                )
            else:
                return self._create_error_result("No data found in XLSX")

        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            return self._create_error_result(f"XLSX error: {str(e)}")

    def _extract_sheet_text(self, worksheet) -> str:
        """
        워크시트에서 텍스트를 추출합니다.

        Args:
            worksheet: openpyxl 워크시트 객체

        Returns:
            str: 추출된 시트 텍스트
        """
        sheet_rows = []

        # 헤더 행 추출 (빈 워크시트 가드)
        if worksheet.max_row is None or worksheet.max_row < 1:
            return ""

        headers = []
        for cell in worksheet[1]:
            headers.append(str(cell.value) if cell.value is not None else "")
        if any(headers):
            sheet_rows.append("Headers: " + ", ".join(headers))

        # 데이터 행 추출 (최대 MAX_ROWS_PER_SHEET개)
        row_count = 0
        for row in worksheet.iter_rows(min_row=2, values_only=False):
            if row_count >= self.MAX_ROWS_PER_SHEET:
                break

            row_values = []
            for col_idx, cell in enumerate(row):
                cell_value = cell.value
                if cell_value is not None:
                    header = headers[col_idx] if col_idx < len(headers) else f'Col{col_idx}'
                    row_values.append(f"{header}: {cell_value}")

            if row_values:
                sheet_rows.append(", ".join(row_values))
                row_count += 1

        return "\n".join(sheet_rows) if sheet_rows else ""
